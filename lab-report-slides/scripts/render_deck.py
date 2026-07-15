#!/usr/bin/env python3
"""Render a structured lab-report deck to HTML, PDF, PNG slides, and image-only PPTX."""

from __future__ import annotations

import argparse
import base64
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Any


SLIDE_WIDTH = 1600
SLIDE_HEIGHT = 900
BASE_CSS = """
@page { size: 16in 9in; margin: 0; }
* { box-sizing: border-box; }
html, body { margin: 0; padding: 0; background: #e8edf3; }
body { font-family: "DengXian", "等线", "Microsoft YaHei", Arial, sans-serif; color: #1f2937; }
.deck { width: 1600px; margin: 0 auto; }
.slide { width: 1600px; height: 900px; position: relative; overflow: hidden; page-break-after: always; background: #fff; border-top: 8px solid #4472c4; padding: 58px 76px 52px; }
.slide:last-child { page-break-after: auto; }
.kicker { color: #5b9bd5; font-size: 20px; line-height: 1.2; font-weight: 700; margin: 0 0 14px; }
h1 { color: #24364b; font-size: 42px; line-height: 1.16; margin: 0; font-weight: 700; letter-spacing: 0; }
.subtitle { color: #6b7280; font-size: 21px; line-height: 1.35; margin-top: 14px; }
.header { min-height: 104px; }
.content { height: 670px; margin-top: 24px; display: flex; gap: 30px; align-items: stretch; }
.content.single { display: block; }
.column { flex: 1; min-width: 0; display: flex; flex-direction: column; gap: 18px; }
.column.wide { flex: 1.55; }
.block { min-width: 0; }
.block h2 { color: #24364b; font-size: 25px; line-height: 1.2; margin: 0 0 8px; }
.block p, .block li { font-size: 23px; line-height: 1.42; margin: 0; }
.block ul { margin: 0; padding-left: 28px; }
.block li + li { margin-top: 9px; }
.callout { border-left: 7px solid #5b9bd5; background: #f3f6fa; padding: 20px 24px; }
.status { position: absolute; right: 76px; top: 64px; color: #4472c4; font-size: 19px; font-weight: 700; }
.metric-row { display: flex; gap: 18px; flex-wrap: wrap; }
.metric { flex: 1 1 190px; border: 2px solid #d9e2f0; padding: 16px 20px; min-height: 110px; }
.metric .value { color: #1f4e79; font-size: 36px; line-height: 1.05; font-weight: 700; }
.metric .label { color: #6b7280; font-size: 18px; margin-top: 8px; }
figure { margin: 0; min-width: 0; height: 100%; display: flex; flex-direction: column; }
figure img { display: block; width: 100%; height: 100%; min-height: 0; object-fit: contain; background: #fbfcfe; border: 1px solid #d9e2f0; }
figcaption { color: #6b7280; font-size: 16px; line-height: 1.25; margin-top: 8px; }
.image-block { flex: 1; min-height: 0; }
.image-block figure { height: 100%; }
.image-block img { max-height: 590px; }
.missing { border: 2px dashed #c96b6b; color: #a33b3b; min-height: 180px; display: grid; place-items: center; padding: 20px; font-size: 20px; text-align: center; }
.footer { position: absolute; left: 76px; right: 76px; bottom: 22px; display: flex; justify-content: space-between; color: #8a94a3; font-size: 15px; }
@media print { html, body { background: #fff; } .slide { box-shadow: none; } }
"""


def esc(value: Any) -> str:
    return html.escape(str(value or ""), quote=True)


def data_uri(path_value: str | None) -> str | None:
    if not path_value:
        return None
    path = Path(path_value).expanduser()
    if not path.exists() or not path.is_file():
        return None
    mime = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    try:
        encoded = base64.b64encode(path.read_bytes()).decode("ascii")
    except OSError:
        return None
    return f"data:{mime};base64,{encoded}"


def lines_html(text: str) -> str:
    rows = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    if not rows:
        return ""
    if len(rows) == 1:
        return f"<p>{esc(rows[0])}</p>"
    return "<ul>" + "".join(f"<li>{esc(row.lstrip('-* '))}</li>" for row in rows) + "</ul>"


def render_block(block: dict[str, Any]) -> str:
    kind = str(block.get("type") or block.get("kind") or "text")
    heading = block.get("heading") or block.get("title")
    if kind == "image":
        src = data_uri(block.get("path"))
        caption = esc(block.get("caption") or block.get("source") or "")
        if not src:
            body = f"<div class=\"missing\">[MISSING: {esc(block.get('path'))}]</div>"
        else:
            body = f"<figure><img src=\"{src}\" alt=\"{caption}\"><figcaption>{caption}</figcaption></figure>"
        return f"<div class=\"block image-block\">{body}</div>"
    if kind == "metric":
        return f"<div class=\"metric\"><div class=\"value\">{esc(block.get('value'))}</div><div class=\"label\">{esc(block.get('label'))}</div></div>"
    body = lines_html(block.get("text") or block.get("body") or "")
    if kind == "callout":
        return f"<div class=\"block callout\">{f'<h2>{esc(heading)}</h2>' if heading else ''}{body}</div>"
    return f"<div class=\"block\">{f'<h2>{esc(heading)}</h2>' if heading else ''}{body}</div>"


def slide_html(slide: dict[str, Any], index: int, total: int, footer: str) -> str:
    title = slide.get("title") or f"Slide {index}"
    kicker = slide.get("kicker") or slide.get("section") or ""
    subtitle = slide.get("subtitle") or ""
    status = slide.get("status") or ""
    blocks = list(slide.get("blocks") or [])
    if not blocks and slide.get("body"):
        blocks = [{"type": "text", "text": slide["body"]}]
    images = [block for block in blocks if str(block.get("type") or block.get("kind")) == "image"]
    non_images = [block for block in blocks if block not in images]
    if images and non_images:
        left = "".join(render_block(block) for block in non_images)
        right = "".join(render_block(block) for block in images)
        content = f'<div class="column wide">{left}</div><div class="column">{right}</div>'
    else:
        content = "".join(render_block(block) for block in blocks)
        content = f'<div class="column">{content}</div>'
    return (
        f'<section class="slide" data-slide="{index}">'
        f'<div class="header"><div class="kicker">{esc(kicker)}</div><h1>{esc(title)}</h1>'
        f'{f"<div class=\"subtitle\">{esc(subtitle)}</div>" if subtitle else ""}</div>'
        f'{f"<div class=\"status\">{esc(status)}</div>" if status else ""}'
        f'<div class="content">{content}</div>'
        f'<div class="footer"><span>{esc(footer)}</span><span>{index} / {total}</span></div>'
        "</section>"
    )


def deck_html(deck: dict[str, Any], slides: list[dict[str, Any]] | None = None) -> str:
    selected = slides if slides is not None else list(deck.get("slides") or [])
    if not selected:
        selected = [{"title": "暂无可汇报内容", "body": "未找到当天的有效工作记录。请检查会话范围或先完成提纲确认。"}]
    footer = deck.get("footer") or deck.get("date") or deck.get("target_date") or "科研工作汇报"
    body = "".join(slide_html(slide, idx, len(selected), footer) for idx, slide in enumerate(selected, 1))
    title = deck.get("title") or "科研工作汇报"
    return f'<!doctype html><html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=1600, initial-scale=1"><title>{esc(title)}</title><style>{BASE_CSS}</style></head><body><main class="deck">{body}</main></body></html>'


def find_browser() -> str:
    candidates = [
        os.environ.get("BROWSER_PATH"),
        shutil.which("msedge"),
        shutil.which("chrome"),
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
        str(Path.home() / r"AppData\Local\Microsoft\Edge\Application\msedge.exe"),
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    raise RuntimeError("未找到 Edge 或 Chrome。请设置 BROWSER_PATH。")


def browser_args(browser: str, profile: Path) -> list[str]:
    return [
        browser,
        "--headless=new",
        "--disable-gpu",
        "--disable-extensions",
        "--no-first-run",
        "--no-default-browser-check",
        "--allow-file-access-from-files",
        "--run-all-compositor-stages-before-draw",
        "--virtual-time-budget=1200",
        f"--user-data-dir={profile}",
    ]


def export_browser(html_path: Path, pdf_path: Path, slide_paths: list[Path], slides: list[dict[str, Any]], deck: dict[str, Any], temp_dir: Path) -> None:
    browser = find_browser()
    profile = temp_dir / "browser-profile"
    uri = html_path.resolve().as_uri()
    subprocess.run(browser_args(browser, profile) + [f"--print-to-pdf={pdf_path}", "--print-to-pdf-no-header", uri], check=True, capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=90)
    for index, target in enumerate(slide_paths, 1):
        single = temp_dir / f"slide-{index}.html"
        single.write_text(deck_html(deck, [slides[index - 1]]), encoding="utf-8")
        subprocess.run(browser_args(browser, profile) + [f"--window-size={SLIDE_WIDTH},{SLIDE_HEIGHT}", f"--screenshot={target}", single.resolve().as_uri()], check=True, capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=60)


def choose_stem(output_dir: Path, requested: str) -> str:
    pattern = re.compile(rf"^{re.escape(requested)}(?:_v(\d+))?$")
    existing: list[int] = []
    for path in output_dir.iterdir() if output_dir.exists() else []:
        match = pattern.match(path.stem)
        if match:
            existing.append(int(match.group(1) or 1))
    if not existing:
        return requested
    return f"{requested}_v{max(existing) + 1}"


def make_pptx(slide_paths: list[Path], output_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx。请安装 python-pptx 后重试。") from exc
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for image_path in slide_paths:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(str(image_path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    presentation.save(output_path)


def render(deck_path: Path, output_dir: Path, requested_stem: str) -> dict[str, Any]:
    deck = json.loads(deck_path.read_text(encoding="utf-8-sig"))
    output_dir.mkdir(parents=True, exist_ok=True)
    stem = choose_stem(output_dir, requested_stem)
    html_path = output_dir / f"{stem}.html"
    pdf_path = output_dir / f"{stem}.pdf"
    pptx_path = output_dir / f"{stem}.pptx"
    manifest_path = output_dir / f"{stem}.manifest.json"
    slides = list(deck.get("slides") or [])
    if not slides:
        slides = [{"title": "暂无可汇报内容", "body": "未找到当天的有效工作记录。"}]
    html_path.write_text(deck_html(deck, slides), encoding="utf-8")
    slide_paths = [output_dir / f"{stem}_{index:02d}.png" for index in range(1, len(slides) + 1)]
    with tempfile.TemporaryDirectory(prefix="lab-report-slides-") as temp:
        export_browser(html_path, pdf_path, slide_paths, slides, deck, Path(temp))
    make_pptx(slide_paths, pptx_path)
    if not pdf_path.exists() or pdf_path.stat().st_size == 0:
        raise RuntimeError("PDF 导出失败或文件为空。")
    if not pptx_path.exists() or not zipfile.is_zipfile(pptx_path):
        raise RuntimeError("PPTX 导出失败或不是有效的 Office 文件。")
    manifest = {
        "schema_version": 1,
        "requested_stem": requested_stem,
        "stem": stem,
        "slide_count": len(slides),
        "files": {"html": str(html_path), "pdf": str(pdf_path), "pptx": str(pptx_path), "png": [str(path) for path in slide_paths]},
    }
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return manifest


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--deck", required=True, help="Structured deck JSON")
    parser.add_argument("--output-dir", required=True)
    parser.add_argument("--base-name", required=True, help="YYYYMMDD or YYYYMMDD组会")
    args = parser.parse_args()
    result = render(Path(args.deck), Path(args.output_dir), args.base_name)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
