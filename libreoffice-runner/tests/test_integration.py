from __future__ import annotations

import base64
import hashlib
import json
import os
import subprocess
import tempfile
import time
import unittest
import zipfile
from pathlib import Path
from xml.etree import ElementTree

from _support import PYTHON, RUNNER_CLI, WindowsOnlyTestCase


class LibreOfficeIntegrationTests(WindowsOnlyTestCase):
    @classmethod
    def setUpClass(cls) -> None:
        super().setUpClass()
        if os.environ.get("RUN_LIBREOFFICE_INTEGRATION") != "1":
            raise unittest.SkipTest("Set RUN_LIBREOFFICE_INTEGRATION=1 to run real LibreOffice tests")
        try:
            import psutil
        except ImportError as exc:
            raise unittest.SkipTest("psutil is required for process-protection assertions") from exc
        cls.psutil = psutil
        existing = cls._libreoffice_processes()
        if existing:
            raise unittest.SkipTest("LibreOffice is already running; integration test will not touch user processes")
        cls.bootstrap = Path(r"C:\Program Files\LibreOffice\program\bootstrap.ini")
        cls.bootstrap_hash = hashlib.sha256(cls.bootstrap.read_bytes()).hexdigest()
        cls.default_profile = Path(os.environ["APPDATA"]) / "LibreOffice" / "4" / "user" / "registrymodifications.xcu"
        cls.default_profile_before = cls._path_hash(cls.default_profile)

    @classmethod
    def _libreoffice_processes(cls) -> list[object]:
        names = {"soffice.exe", "soffice.com", "soffice.bin"}
        return [
            process
            for process in cls.psutil.process_iter(["name"])
            if (process.info.get("name") or "").lower() in names
        ]

    def setUp(self) -> None:
        self.temp_dir = tempfile.TemporaryDirectory(prefix="lo-runner-integration-")
        self.root = Path(self.temp_dir.name)
        self.inputs = self.root / "inputs"
        self.outputs = self.root / "outputs"
        self.inputs.mkdir()
        self.outputs.mkdir()

    def tearDown(self) -> None:
        self.temp_dir.cleanup()
        self.assertEqual(hashlib.sha256(self.bootstrap.read_bytes()).hexdigest(), self.bootstrap_hash)
        self.assertEqual(self._path_hash(self.default_profile), self.default_profile_before)
        self.assertEqual(self._libreoffice_processes(), [])

    def test_mixed_parallel_conversions_and_accept_changes(self) -> None:
        docx_a = self.inputs / "first" / "same-name.docx"
        docx_b = self.inputs / "second" / "same-name.docx"
        tracked = self.inputs / "tracked.docx"
        workbook = self.inputs / "formula.xlsx"
        pptx_a = self.inputs / "slides-a.pptx"
        pptx_b = self.inputs / "slides-b.pptx"
        docx_a.parent.mkdir()
        docx_b.parent.mkdir()
        self._build_docx(docx_a, "Alpha")
        self._build_docx(docx_b, "Beta")
        self._build_tracked_docx(tracked)
        self._build_xlsx(workbook)
        self._build_pptx(pptx_a, "A")
        self._build_pptx(pptx_b, "B")

        jobs = [
            ("pdf", docx_a, self.outputs / "alpha.pdf"),
            ("pdf", docx_b, self.outputs / "beta.pdf"),
            ("recalc", workbook, self.outputs / "formula-recalc.xlsx"),
            ("pdf", workbook, self.outputs / "formula.pdf"),
            ("pdf", pptx_a, self.outputs / "slides-a.pdf"),
            ("pdf", pptx_b, self.outputs / "slides-b.pdf"),
        ]
        processes = [
            subprocess.Popen(
                [str(PYTHON), str(RUNNER_CLI), operation, str(source), str(output), "--run-timeout", "120"],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace",
            )
            for operation, source, output in jobs
        ]
        peak_soffice_bin = 0
        while any(process.poll() is None for process in processes):
            running = [
                process
                for process in self._libreoffice_processes()
                if (process.info.get("name") or "").lower() == "soffice.bin"
            ]
            peak_soffice_bin = max(peak_soffice_bin, len(running))
            time.sleep(0.05)
        reports = []
        for process in processes:
            stdout, stderr = process.communicate(timeout=10)
            self.assertEqual(process.returncode, 0, f"stdout={stdout}\nstderr={stderr}")
            reports.append(json.loads(stdout))
        self.assertLessEqual(peak_soffice_bin, 2)
        self.assertEqual(len(reports), 6)
        self.assertTrue(all(report["ok"] for report in reports))
        for report in reports:
            profile_argument = next(
                argument for argument in report["command"] if argument.startswith("-env:UserInstallation=")
            )
            self.assertIn("file:///", profile_argument)
            self.assertIn("sanan-lo-", profile_argument)
        self.assertIn(b"/Image", (self.outputs / "alpha.pdf").read_bytes())
        self.assertEqual(self._formula_cache(self.outputs / "formula-recalc.xlsx"), "5")
        self.assertEqual(reports[4]["validation"]["pages"], 2)
        self.assertEqual(reports[5]["validation"]["pages"], 2)

        result = subprocess.run(
            [str(PYTHON), str(RUNNER_CLI), "accept-changes", str(tracked), str(self.outputs / "accepted.docx")],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=180,
            check=False,
        )
        self.assertEqual(result.returncode, 0, f"stdout={result.stdout}\nstderr={result.stderr}")
        report = json.loads(result.stdout)
        self.assertTrue(report["ok"])
        self.assertEqual(Path(report["command"][0]).name.lower(), "soffice.com")
        accept_profile = next(
            argument for argument in report["command"] if argument.startswith("-env:UserInstallation=")
        )
        self.assertIn("sanan-lo-", accept_profile)
        with zipfile.ZipFile(self.outputs / "accepted.docx") as archive:
            document = archive.read("word/document.xml")
        self.assertNotIn(b"<w:ins", document)
        self.assertNotIn(b"<w:del", document)

    @staticmethod
    def _path_hash(path: Path) -> tuple[bool, str | None]:
        if not path.exists():
            return False, None
        return True, hashlib.sha256(path.read_bytes()).hexdigest()

    def _build_docx(self, path: Path, text: str) -> None:
        from docx import Document

        image = path.with_suffix(".png")
        image.write_bytes(
            base64.b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9Y9jY5kAAAAASUVORK5CYII="
            )
        )
        document = Document()
        document.add_paragraph(text)
        document.add_picture(str(image))
        document.save(path)

    def _build_tracked_docx(self, path: Path) -> None:
        from docx import Document

        document = Document()
        document.add_paragraph("Before")
        document.save(path)
        replacement = (
            b'<w:ins w:id="1" w:author="Runner" w:date="2026-07-23T00:00:00Z">'
            b"<w:r><w:t>Inserted</w:t></w:r></w:ins>"
            b'<w:del w:id="2" w:author="Runner" w:date="2026-07-23T00:00:00Z">'
            b"<w:r><w:delText>Deleted</w:delText></w:r></w:del>"
        )
        with zipfile.ZipFile(path, "r") as source:
            parts = {name: source.read(name) for name in source.namelist()}
        parts["word/document.xml"] = parts["word/document.xml"].replace(b"</w:p>", replacement + b"</w:p>", 1)
        temporary = path.with_suffix(".tmp")
        with zipfile.ZipFile(temporary, "w", zipfile.ZIP_DEFLATED) as target:
            for name, payload in parts.items():
                target.writestr(name, payload)
        temporary.replace(path)

    def _build_xlsx(self, path: Path) -> None:
        from openpyxl import Workbook

        workbook = Workbook()
        sheet = workbook.active
        sheet["A1"] = 2
        sheet["B1"] = 3
        sheet["C1"] = "=SUM(A1:B1)"
        workbook.save(path)

    def _build_pptx(self, path: Path, text: str) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = Inches(7.5)
        for index in range(2):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = f"{text}-{index + 1}"
        presentation.save(path)

    @staticmethod
    def _formula_cache(path: Path) -> str | None:
        namespace = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
        with zipfile.ZipFile(path) as archive:
            root = ElementTree.fromstring(archive.read("xl/worksheets/sheet1.xml"))
        for cell in root.findall(f".//{namespace}c"):
            if cell.attrib.get("r") == "C1":
                value = cell.find(f"{namespace}v")
                return value.text if value is not None else None
        return None
